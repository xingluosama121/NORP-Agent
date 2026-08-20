# ──────────────────────────────────────────────────────────────
# Plugin: Document Reader (文档阅读器)
# Publisher: xingluosama
# Version: 1.0.0
# Description: 让 Agent 读取 Office 文档：Word (.docx)、
#   Excel (.xlsx)、PowerPoint (.pptx)。自动提取文字、
#   表格、幻灯片内容为结构化 Markdown。
#
# 依赖（需预先安装）：
#   pip install python-docx openpyxl python-pptx
#
# 若未安装对应库，工具会返回清晰的安装提示。
# ──────────────────────────────────────────────────────────────

PLUGIN_NAME = "Document Reader"
PLUGIN_PUBLISHER = "xingluosama"
PLUGIN_VERSION = "1.0.0"
PLUGIN_DESCRIPTION = (
    "文档阅读器：读取 .docx / .xlsx / .pptx 文件，"
    "提取文字、表格、幻灯片内容为结构化 Markdown。"
)

import os
import re
from typing import Any, Dict, List, Optional, Tuple

# ═══════════════════════════════════════════════════════════════
#  工具注册
# ═══════════════════════════════════════════════════════════════

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "read_docx",
            "description": (
                "读取 Word 文档 (.docx)，提取文本内容为 Markdown 格式。"
                "支持提取段落文字、表格、页眉页脚。"
                "适合阅读报告、合同、论文、手册等 Word 文档。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": (
                            "文档路径。可以是相对于当前工作区的路径，也可以是绝对路径。"
                            "例如 'docs/report.docx' 或 '/tmp/contract.docx'。"
                        )
                    },
                    "include_tables": {
                        "type": "boolean",
                        "description": "是否提取表格内容，默认 true。表格会渲染为 Markdown 表格。",
                        "default": True
                    },
                    "include_headers_footers": {
                        "type": "boolean",
                        "description": "是否提取页眉页脚，默认 false。",
                        "default": False
                    },
                    "max_paragraphs": {
                        "type": "integer",
                        "description": "最大返回段落数，默认 0 表示不限制。用于控制输出长度。",
                        "default": 0
                    }
                },
                "required": ["file_path"],
                "additionalProperties": False
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_xlsx",
            "description": (
                "读取 Excel 电子表格 (.xlsx)，提取数据为 Markdown 表格。"
                "支持选择工作表、列范围、行范围，以及将首行作为表头。"
                "适合阅读数据报表、统计表、清单等 Excel 文件。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": (
                            "文档路径。可以是相对于当前工作区的路径，也可以是绝对路径。"
                            "例如 'data/finance.xlsx'。"
                        )
                    },
                    "sheet_name": {
                        "type": "string",
                        "description": (
                            "工作表名称。不指定则读取第一个工作表。"
                            "可先用 sheet_name='__list__' 列出所有工作表名称。"
                        )
                    },
                    "cell_range": {
                        "type": "string",
                        "description": (
                            "单元格范围，如 'A1:F100'。不指定则读取整个工作表。"
                            "超出实际数据范围的区域会被自动裁剪。"
                        )
                    },
                    "max_rows": {
                        "type": "integer",
                        "description": "最大返回行数（不含表头），默认 200。",
                        "default": 200
                    },
                    "max_cols": {
                        "type": "integer",
                        "description": "最大返回列数，默认 50。",
                        "default": 50
                    },
                    "header_row": {
                        "type": "integer",
                        "description": "表头所在行号（从 1 开始），默认 1。设为 0 则无表头。",
                        "default": 1
                    }
                },
                "required": ["file_path"],
                "additionalProperties": False
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_pptx",
            "description": (
                "读取 PowerPoint 演示文稿 (.pptx)，提取幻灯片内容为 Markdown。"
                "每张幻灯片提取标题、正文、备注。"
                "适合阅读演示文稿、培训材料、提案等 PPT 文件。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": (
                            "文档路径。可以是相对于当前工作区的路径，也可以是绝对路径。"
                            "例如 'slides/presentation.pptx'。"
                        )
                    },
                    "slide_range": {
                        "type": "string",
                        "description": (
                            "幻灯片范围，如 '1'（第 1 张）、'1-5'（第 1 到 5 张）、"
                            "'1,3,5'（第 1、3、5 张）。不指定则读取全部幻灯片。"
                        )
                    },
                    "include_notes": {
                        "type": "boolean",
                        "description": "是否提取演讲者备注，默认 true。",
                        "default": True
                    },
                    "include_shapes": {
                        "type": "boolean",
                        "description": (
                            "是否提取形状中的文本（如文本框、图形中的文字），默认 true。"
                        ),
                        "default": True
                    },
                    "max_slides": {
                        "type": "integer",
                        "description": "最大返回幻灯片数，默认 0 表示不限制。",
                        "default": 0
                    }
                },
                "required": ["file_path"],
                "additionalProperties": False
            }
        }
    }
]

# ═══════════════════════════════════════════════════════════════
#  Lazy imports – 优雅处理依赖缺失
# ═══════════════════════════════════════════════════════════════

_MISSING_DOCX: Optional[str] = None
_MISSING_XLSX: Optional[str] = None
_MISSING_PPTX: Optional[str] = None

def _get_docx_module():
    """Lazy-import python-docx with a helpful error message."""
    global _MISSING_DOCX
    if _MISSING_DOCX:
        return None, _MISSING_DOCX
    try:
        import docx
        return docx, None
    except ImportError:
        _MISSING_DOCX = (
            "❌ 未安装 python-docx 库。\n"
            "请运行：`pip install python-docx`\n"
            "然后重试。"
        )
        return None, _MISSING_DOCX

def _get_xlsx_module():
    """Lazy-import openpyxl with a helpful error message."""
    global _MISSING_XLSX
    if _MISSING_XLSX:
        return None, _MISSING_XLSX
    try:
        import openpyxl
        return openpyxl, None
    except ImportError:
        _MISSING_XLSX = (
            "❌ 未安装 openpyxl 库。\n"
            "请运行：`pip install openpyxl`\n"
            "然后重试。"
        )
        return None, _MISSING_XLSX

def _get_pptx_module():
    """Lazy-import python-pptx with a helpful error message."""
    global _MISSING_PPTX
    if _MISSING_PPTX:
        return None, _MISSING_PPTX
    try:
        import pptx
        return pptx, None
    except ImportError:
        _MISSING_PPTX = (
            "❌ 未安装 python-pptx 库。\n"
            "请运行：`pip install python-pptx`\n"
            "然后重试。"
        )
        return None, _MISSING_PPTX

# ═══════════════════════════════════════════════════════════════
#  Helpers
# ═══════════════════════════════════════════════════════════════

def _resolve_path(file_path: str, context) -> str:
    """Resolve file path relative to workspace root or as absolute."""
    if os.path.isabs(file_path):
        return file_path
    project_root = getattr(context, "project_root", "") or os.getcwd()
    return os.path.normpath(os.path.join(project_root, file_path))


def _validate_file(file_path: str, expected_ext: str) -> Optional[str]:
    """Validate that file exists and has the correct extension.

    Returns an error string or None if valid.
    """
    if not os.path.exists(file_path):
        return f"❌ 文件不存在：`{file_path}`"
    if not file_path.lower().endswith(expected_ext):
        return (
            f"❌ 文件扩展名不匹配：期望 `{expected_ext}`，"
            f"实际为 `{os.path.splitext(file_path)[1]}`"
        )
    return None


def _safe_text(text) -> str:
    """Safely convert any text to string, handling None and non-strings."""
    if text is None:
        return ""
    return str(text).strip()


def _limit_output(lines: List[str], max_items: int, label: str) -> List[str]:
    """Truncate output lines with a notice."""
    if max_items > 0 and len(lines) > max_items:
        truncated = lines[:max_items]
        truncated.append(
            f"\n> ⚠️ 输出被截断：仅显示前 {max_items} 项，"
            f"共 {len(lines)} 项。请缩小范围重新查询。"
        )
        return truncated
    return lines

# ═══════════════════════════════════════════════════════════════
#  DOCX Reader
# ═══════════════════════════════════════════════════════════════

def _read_docx(file_path: str,
               include_tables: bool = True,
               include_headers_footers: bool = False,
               max_paragraphs: int = 0) -> str:
    """Read a .docx file and return Markdown text."""

    docx, err = _get_docx_module()
    if err:
        return err

    error = _validate_file(file_path, ".docx")
    if error:
        return error

    try:
        doc = docx.Document(file_path)
    except Exception as exc:
        return f"❌ 无法打开文档 `{file_path}`：{exc}"

    output: List[str] = []

    # ── 页眉页脚 ──
    if include_headers_footers:
        header_texts: List[str] = []
        for section in doc.sections:
            if section.header and section.header.paragraphs:
                header_texts.extend(
                    p.text for p in section.header.paragraphs if p.text.strip()
                )
            if section.footer and section.footer.paragraphs:
                header_texts.extend(
                    p.text for p in section.footer.paragraphs if p.text.strip()
                )
        if header_texts:
            output.append("---")
            output.append("### 📋 页眉/页脚")
            output.append("")
            for ht in header_texts:
                output.append(f"> {ht}")
            output.append("---")
            output.append("")

    # ── 计数 ──
    all_paras = doc.paragraphs
    all_tables = doc.tables
    output.append(
        f"📄 **文档信息**：{len(all_paras)} 个段落，{len(all_tables)} 个表格\n"
    )

    # ── 遍历 body 元素（保持段落与表格的顺序）──
    body = doc.element.body
    para_idx = 0
    table_idx = 0
    para_count = 0
    # 用 xml 标签区分段落和表格
    from docx.oxml.ns import qn

    for child in body:
        tag = child.tag

        # 段落
        if tag == qn("w:p"):
            if para_idx >= len(all_paras):
                para_idx += 1
                continue
            para = all_paras[para_idx]
            para_idx += 1

            text = para.text.strip()
            if not text:
                continue

            # 检测标题样式
            style_name = (para.style.name if para.style else "").lower()
            prefix = ""
            if "heading 1" in style_name or style_name == "title":
                prefix = "# "
            elif "heading 2" in style_name:
                prefix = "## "
            elif "heading 3" in style_name:
                prefix = "### "
            elif "heading 4" in style_name:
                prefix = "#### "

            output.append(f"{prefix}{text}")
            output.append("")
            para_count += 1

        # 表格
        elif tag == qn("w:tbl"):
            if not include_tables or table_idx >= len(all_tables):
                table_idx += 1
                continue
            table = all_tables[table_idx]
            table_idx += 1

            table_md = _table_to_markdown(table)
            if table_md:
                output.append(table_md)
                output.append("")
                para_count += 1

    # 如果 docx 库版本不同导致 body 元素枚举为空，回退到单独枚举
    if para_count == 0:
        for para in all_paras:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name if para.style else "").lower()
            prefix = ""
            if "heading 1" in style_name or style_name == "title":
                prefix = "# "
            elif "heading 2" in style_name:
                prefix = "## "
            elif "heading 3" in style_name:
                prefix = "### "
            elif "heading 4" in style_name:
                prefix = "#### "
            output.append(f"{prefix}{text}")
            output.append("")
            para_count += 1

        if include_tables:
            for table in all_tables:
                table_md = _table_to_markdown(table)
                if table_md:
                    output.append(table_md)
                    output.append("")
                    para_count += 1

    output = _limit_output(output, max_paragraphs, "段落/表格") if max_paragraphs > 0 else output

    if not output or (len(output) <= 2 and para_count == 0):
        return "📄 文档为空或无法提取文本内容。"

    return "\n".join(output)


def _table_to_markdown(table) -> str:
    """Convert a docx Table to Markdown table."""
    rows = table.rows
    if not rows:
        return ""

    max_cols = max(len(row.cells) for row in rows)
    if max_cols == 0:
        return ""

    md_rows: List[str] = []

    for r_idx, row in enumerate(rows):
        cells = []
        for c_idx in range(max_cols):
            if c_idx < len(row.cells):
                cell_text = row.cells[c_idx].text.replace("\n", " ").replace("|", "\\|")
                cells.append(cell_text.strip())
            else:
                cells.append("")
        md_rows.append("| " + " | ".join(cells) + " |")

        # 表头分隔线（第一行后）
        if r_idx == 0:
            md_rows.append("| " + " | ".join(["---"] * max_cols) + " |")

    return "\n".join(md_rows)


# ═══════════════════════════════════════════════════════════════
#  XLSX Reader
# ═══════════════════════════════════════════════════════════════

def _parse_cell_range(cell_range: str) -> Tuple[int, int, int, int]:
    """Parse 'A1:F100' into (min_col, min_row, max_col, max_row). 1-indexed.

    Returns (-1, -1, -1, -1) for invalid ranges.
    """
    if not cell_range or ":" not in cell_range:
        # Try single cell like 'A1'
        match = re.match(r"^([A-Z]+)(\d+)$", cell_range.strip().upper())
        if match:
            col = _col_letter_to_num(match.group(1))
            row = int(match.group(2))
            return (col, row, col, row)
        return (-1, -1, -1, -1)

    parts = cell_range.split(":", 1)
    start_match = re.match(r"^([A-Z]+)(\d+)$", parts[0].strip().upper())
    end_match = re.match(r"^([A-Z]+)(\d+)$", parts[1].strip().upper())

    if not start_match or not end_match:
        return (-1, -1, -1, -1)

    c1 = _col_letter_to_num(start_match.group(1))
    r1 = int(start_match.group(2))
    c2 = _col_letter_to_num(end_match.group(1))
    r2 = int(end_match.group(2))

    min_c, max_c = min(c1, c2), max(c1, c2)
    min_r, max_r = min(r1, r2), max(r1, r2)
    return (min_c, min_r, max_c, max_r)


def _col_letter_to_num(letter: str) -> int:
    """Convert 'A' → 1, 'Z' → 26, 'AA' → 27, etc."""
    num = 0
    for ch in letter.upper():
        num = num * 26 + (ord(ch) - ord("A") + 1)
    return num


def _col_num_to_letter(num: int) -> str:
    """Convert 1 → 'A', 26 → 'Z', 27 → 'AA'."""
    result = ""
    while num > 0:
        num, rem = divmod(num - 1, 26)
        result = chr(ord("A") + rem) + result
    return result


def _read_xlsx(file_path: str,
               sheet_name: str = "",
               cell_range: str = "",
               max_rows: int = 200,
               max_cols: int = 50,
               header_row: int = 1) -> str:
    """Read an .xlsx file and return Markdown tables."""

    openpyxl, err = _get_xlsx_module()
    if err:
        return err

    error = _validate_file(file_path, ".xlsx")
    if error:
        return error

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    except Exception as exc:
        return f"❌ 无法打开电子表格 `{file_path}`：{exc}"

    # ── 列出工作表 ──
    if sheet_name == "__list__":
        sheets = wb.sheetnames
        lines = [f"📊 **工作表列表** ({len(sheets)} 个)：", ""]
        for i, name in enumerate(sheets, 1):
            ws = wb[name]
            lines.append(
                f"  **{i}.** `{name}` — "
                f"{ws.max_row} 行 × {ws.max_column} 列"
            )
        wb.close()
        return "\n".join(lines)

    # ── 选择工作表 ──
    if sheet_name:
        if sheet_name not in wb.sheetnames:
            wb.close()
            similar = [s for s in wb.sheetnames if sheet_name.lower() in s.lower()]
            hint = f" 可用的工作表：{', '.join(wb.sheetnames[:10])}"
            if similar:
                hint += f"\n  相似的工作表：{', '.join(similar)}"
            return f"❌ 工作表 `{sheet_name}` 不存在。{hint}"
        ws = wb[sheet_name]
    else:
        ws = wb.active

    # ── 确定数据范围 ──
    data_min_row = 1
    data_max_row = ws.max_row or 1
    data_min_col = 1
    data_max_col = ws.max_column or 1

    if cell_range:
        c1, r1, c2, r2 = _parse_cell_range(cell_range)
        if c1 > 0:
            data_min_col = max(c1, 1)
            data_max_col = min(c2, data_max_col)
            data_min_row = max(r1, 1)
            data_max_row = min(r2, data_max_row)

    # 限制行列数
    data_max_row = min(data_max_row, data_min_row + max_rows - 1)
    data_max_col = min(data_max_col, data_min_col + max_cols - 1)

    # ── 读取数据 ──
    rows_data: List[List[str]] = []
    last_data_row = data_min_row

    for row_idx in range(data_min_row, data_max_row + 1):
        row_values: List[str] = []
        has_any = False
        for col_idx in range(data_min_col, data_max_col + 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            val = cell.value
            if val is not None:
                has_any = True
                # 处理各种类型
                if isinstance(val, float):
                    if val == int(val) and abs(val) < 1e12:
                        row_values.append(str(int(val)))
                    else:
                        row_values.append(str(val))
                else:
                    row_values.append(str(val).replace("\n", " ").replace("|", "\\|").strip())
            else:
                row_values.append("")
        rows_data.append(row_values)
        if has_any:
            last_data_row = row_idx

    wb.close()

    # 裁剪尾部全空行
    while rows_data and all(v == "" for v in rows_data[-1]):
        rows_data.pop()

    if not rows_data:
        return "📊 电子表格为空或指定范围内无数据。"

    # ── 确定表头 ──
    header_labels: List[str] = []
    data_start = 0
    if 0 < header_row <= len(rows_data):
        # header_row 是相对于原始数据的偏移
        hdr_offset = header_row - 1  # 0-indexed within rows_data
        if 0 <= hdr_offset < len(rows_data):
            header_labels = rows_data[hdr_offset]
            data_start = hdr_offset + 1

    # 如果没有表头或全部为空，自动生成列标签
    max_col_count = max((len(r) for r in rows_data), default=0)
    if not header_labels or all(h == "" for h in header_labels):
        header_labels = [_col_num_to_letter(i) for i in range(1, max_col_count + 1)]

    # Ensure all rows have the same width as header
    for row in rows_data:
        while len(row) < len(header_labels):
            row.append("")

    # ── 生成 Markdown 表格 ──
    lines: List[str] = []
    sheet_label = f"`{ws.title}`" if sheet_name else f"`{wb.sheetnames[0]}`"
    lines.append(
        f"📊 **{sheet_label}** — "
        f"{len(rows_data) - data_start} 行 × {len(header_labels)} 列"
    )
    lines.append("")

    # 表头
    lines.append("| " + " | ".join(header_labels) + " |")
    lines.append("| " + " | ".join(["---"] * len(header_labels)) + " |")

    # 数据
    row_count = 0
    for row in rows_data[data_start:]:
        # Pad row to header width
        padded = row + [""] * (len(header_labels) - len(row))
        lines.append("| " + " | ".join(padded[:len(header_labels)]) + " |")
        row_count += 1

    if row_count == 0 and data_start < len(rows_data):
        # No data rows found (header_row consumed all rows)
        lines.append("| " + " | ".join(["(无数据)"] * len(header_labels)) + " |")

    actual_rows = len(rows_data) - data_start
    if actual_rows >= max_rows:
        lines.append(
            f"\n> ⚠️ 表格可能被截断：仅显示前 {max_rows} 行。"
            f"请使用更小的 `cell_range` 或调整 `max_rows`。"
        )

    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════
#  PPTX Reader
# ═══════════════════════════════════════════════════════════════

def _parse_slide_range(slide_range: str, total_slides: int) -> List[int]:
    """Parse '1-5' or '1,3,5' into a list of slide numbers (1-indexed)."""
    if not slide_range or not slide_range.strip():
        return list(range(1, total_slides + 1))

    result: List[int] = []
    parts = slide_range.split(",")
    for part in parts:
        part = part.strip()
        if "-" in part:
            try:
                a, b = part.split("-", 1)
                start = int(a.strip())
                end = int(b.strip())
                result.extend(range(start, end + 1))
            except ValueError:
                pass
        else:
            try:
                result.append(int(part))
            except ValueError:
                pass

    # Filter valid slide numbers
    return [n for n in result if 1 <= n <= total_slides]


def _extract_shape_text(shape) -> str:
    """Extract text from any shape that has a text frame."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _read_pptx(file_path: str,
               slide_range: str = "",
               include_notes: bool = True,
               include_shapes: bool = True,
               max_slides: int = 0) -> str:
    """Read a .pptx file and return Markdown slide content."""

    pptx, err = _get_pptx_module()
    if err:
        return err

    error = _validate_file(file_path, ".pptx")
    if error:
        return error

    try:
        prs = pptx.Presentation(file_path)
    except Exception as exc:
        return f"❌ 无法打开演示文稿 `{file_path}`：{exc}"

    total_slides = len(prs.slides)
    target_slides = _parse_slide_range(slide_range, total_slides)

    if max_slides > 0 and len(target_slides) > max_slides:
        target_slides = target_slides[:max_slides]

    output: List[str] = []
    slide_count = 0

    for slide_num in target_slides:
        slide = prs.slides[slide_num - 1]  # 0-indexed
        slide_count += 1

        output.append(f"## 📽️ 幻灯片 {slide_num} / {total_slides}")
        output.append("")

        # ── 标题 ──
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
        if title_text:
            output.append(f"### {title_text}")
            output.append("")

        # ── 正文占位符 ──
        body_texts: List[str] = []
        for shape in slide.placeholders:
            if shape == slide.shapes.title:
                continue  # 已处理标题
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != title_text:
                    body_texts.append(text)

        # ── 其他形状 ──
        shape_texts: List[str] = []
        if include_shapes:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    continue  # 已在占位符中处理
                text = _extract_shape_text(shape)
                if text and text not in body_texts and text != title_text:
                    shape_texts.append(text)

        # 输出正文
        for bt in body_texts:
            for line in bt.split("\n"):
                line = line.strip()
                if line:
                    output.append(f"- {line}")
            output.append("")

        # 输出形状文本
        if shape_texts:
            for st in shape_texts:
                for line in st.split("\n"):
                    line = line.strip()
                    if line:
                        output.append(f"  - 📌 {line}")
            output.append("")

        # ── 演讲者备注 ──
        if include_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                output.append("> 🎤 **演讲者备注：**")
                for line in notes_text.split("\n"):
                    output.append(f"> {line}")
                output.append("")

        # 如果没有提取到任何内容
        if not title_text and not body_texts and not shape_texts:
            output.append("*(此幻灯片无可提取文本)*")
            output.append("")

        output.append("---")
        output.append("")

    if max_slides > 0 and total_slides > max_slides:
        output.insert(
            0,
            f"> ⚠️ 仅显示前 {max_slides} 张幻灯片，共 {total_slides} 张。"
            f"请使用 `slide_range` 缩小范围。\n"
        )

    if slide_count == 0:
        return "📽️ 演示文稿中没有匹配的幻灯片。"

    output.insert(0, f"📽️ **演示文稿** — {total_slides} 张幻灯片，显示 {slide_count} 张\n")

    return "\n".join(output)


# ═══════════════════════════════════════════════════════════════
#  Hook: 文件变更自动检测
# ═══════════════════════════════════════════════════════════════

def after_tool_call(tool_name: str, args: dict, result: str, context):
    """当其他工具创建了 Office 文档时，提示可以读取。

    检测文件写入操作（write_file、exec_cmd 等）创建了 .docx/.xlsx/.pptx，
    自动附加一条提示。
    """
    if not result:
        return result

    office_exts = (".docx", ".xlsx", ".pptx")

    # 检测 write_file 创建了 Office 文件
    if tool_name == "write_file":
        file_path = args.get("path", "")
        if file_path and any(file_path.lower().endswith(ext) for ext in office_exts):
            ext_label = os.path.splitext(file_path)[1].upper()
            reader_tool = {
                ".docx": "read_docx",
                ".xlsx": "read_xlsx",
                ".pptx": "read_pptx",
            }.get(os.path.splitext(file_path)[1].lower(), "read_docx")

            # Append a subtle hint to the result
            if "✅" in result or "成功" in result or "写入" in result:
                hint = (
                    f"\n\n💡 已创建 {ext_label} 文件。"
                    f"可以调用 `{reader_tool}` 工具读取其内容。"
                )
                # Don't modify the return value; just log it
                context.logger.debug(f"Office file created: {file_path}")

    return result  # Never modify the tool output


def on_agent_init(context):
    """启动时检查依赖可用性。"""
    status: List[str] = []

    docx_mod, _ = _get_docx_module()
    xlsx_mod, _ = _get_xlsx_module()
    pptx_mod, _ = _get_pptx_module()

    status.append("✅ DOCX" if docx_mod else "⚠️ DOCX（需 python-docx）")
    status.append("✅ XLSX" if xlsx_mod else "⚠️ XLSX（需 openpyxl）")
    status.append("✅ PPTX" if pptx_mod else "⚠️ PPTX（需 python-pptx）")

    context.logger.info(
        f"Document Reader loaded — {', '.join(status)}"
    )

    # Store dep status for quick checks
    context.storage["doc_reader_deps"] = {
        "docx": docx_mod is not None,
        "xlsx": xlsx_mod is not None,
        "pptx": pptx_mod is not None,
    }


# ═══════════════════════════════════════════════════════════════
#  工具分发
# ═══════════════════════════════════════════════════════════════

def execute(tool_name: str, args: dict, context) -> str:
    """Dispatch tool calls to the appropriate handler."""

    if tool_name == "read_docx":
        file_path = _resolve_path(args.get("file_path", ""), context)
        if not file_path:
            return "❌ 请提供 file_path 参数。"
        return _read_docx(
            file_path=file_path,
            include_tables=args.get("include_tables", True),
            include_headers_footers=args.get("include_headers_footers", False),
            max_paragraphs=args.get("max_paragraphs", 0),
        )

    if tool_name == "read_xlsx":
        file_path = _resolve_path(args.get("file_path", ""), context)
        if not file_path:
            return "❌ 请提供 file_path 参数。"
        return _read_xlsx(
            file_path=file_path,
            sheet_name=args.get("sheet_name", ""),
            cell_range=args.get("cell_range", ""),
            max_rows=args.get("max_rows", 200),
            max_cols=args.get("max_cols", 50),
            header_row=args.get("header_row", 1),
        )

    if tool_name == "read_pptx":
        file_path = _resolve_path(args.get("file_path", ""), context)
        if not file_path:
            return "❌ 请提供 file_path 参数。"
        return _read_pptx(
            file_path=file_path,
            slide_range=args.get("slide_range", ""),
            include_notes=args.get("include_notes", True),
            include_shapes=args.get("include_shapes", True),
            max_slides=args.get("max_slides", 0),
        )

    return f"Unknown tool: {tool_name}"
