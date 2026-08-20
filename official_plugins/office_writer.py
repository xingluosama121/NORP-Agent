# ──────────────────────────────────────────────────────────────
# Plugin: Office Writer (Office 文档写入器)
# Publisher: xingluosama
# Version: 1.0.0
# Description: 让 Agent 创建/写入 Office 文档：
#   Word (.docx)、Excel (.xlsx)、PowerPoint (.pptx)。
#   - write_docx : Markdown 语法 → Word 文档（标题/列表/表格/粗体/斜体/引用）
#   - write_xlsx : 结构化数据 → Excel 工作簿（多工作表/表头/样式/自动列宽）
#   - write_pptx : 结构化数据 → PowerPoint 演示文稿（标题页/要点/备注）
#
# 依赖（需预先安装）：
#   pip install python-docx openpyxl python-pptx
#
# 若未安装对应库，工具会返回清晰的安装提示。
# ──────────────────────────────────────────────────────────────

PLUGIN_NAME = "Office Writer"
PLUGIN_PUBLISHER = "xingluosama"
PLUGIN_VERSION = "1.0.0"
PLUGIN_DESCRIPTION = (
    "Office 文档写入器：创建/写入 .docx / .xlsx / .pptx 文件。"
    "支持 Markdown 转 Word、结构化数据转 Excel、幻灯片转 PowerPoint。"
)

import json
import os
import re
from typing import Any, Dict, List, Optional, Tuple

# ═══════════════════════════════════════════════════════════════
#  工具注册
# ═══════════════════════════════════════════════════════════════

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "write_docx",
            "description": (
                "创建或追加写入 Word 文档 (.docx)。"
                "content 参数使用 Markdown 语法，支持：\n"
                "  - 标题：`# 一级标题` ～ `###### 六级标题`\n"
                "  - 列表：`- 项目`（无序）、`1. 项目`（有序）\n"
                "  - 表格：`| 列1 | 列2 |` 多行，第二行用 `|---|---|` 作分隔\n"
                "  - 引用：`> 引用内容`\n"
                "  - 行内格式：`**粗体**`、`*斜体*`、`` `代码` ``\n"
                "  - 分隔线：`---`\n"
                "适合生成报告、合同、说明书、周报等文档。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": (
                            "输出文件路径。可以是相对于当前工作区的路径，"
                            "也可以是绝对路径。例如 'docs/report.docx'。"
                            "扩展名必须为 .docx，父目录不存在时会自动创建。"
                        )
                    },
                    "content": {
                        "type": "string",
                        "description": (
                            "文档正文，使用 Markdown 语法（标题/列表/表格/引用/粗体/斜体/代码）。"
                            "段落之间用空行分隔。"
                        )
                    },
                    "title": {
                        "type": "string",
                        "description": "文档大标题（显示为文档 Title 样式）。新建文档时生效。",
                        "default": ""
                    },
                    "mode": {
                        "type": "string",
                        "enum": ["create", "append"],
                        "description": (
                            "create：新建文档（默认）；"
                            "append：在已有文档末尾追加内容（文件不存在时自动创建）。"
                        ),
                        "default": "create"
                    },
                    "overwrite": {
                        "type": "boolean",
                        "description": "create 模式下，若文件已存在是否覆盖。默认 true；设为 false 则报错不覆盖。",
                        "default": True
                    }
                },
                "required": ["file_path", "content"],
                "additionalProperties": False
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "write_xlsx",
            "description": (
                "创建 Excel 电子表格 (.xlsx)，支持多个工作表。"
                "每个工作表可指定名称、表头和数据行，表头自动加粗并冻结首行，"
                "列宽按内容自动调整。"
                "适合生成数据报表、统计表、清单、台账等。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": (
                            "输出文件路径。例如 'data/finance.xlsx'。"
                            "扩展名必须为 .xlsx，父目录不存在时会自动创建。"
                        )
                    },
                    "sheets": {
                        "type": "array",
                        "description": (
                            "工作表列表，每个元素格式：\n"
                            "{\"name\": \"工作表名\", \"headers\": [\"列1\", \"列2\"], "
                            "\"rows\": [[\"值\", 123], ...], \"start_cell\": \"A1\"}\n"
                            "name 缺省时自动命名为 Sheet1/Sheet2…；"
                            "headers 缺省时无表头；start_cell 指定数据起始单元格（默认 A1）。"
                        ),
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string", "description": "工作表名称（自动去重）"},
                                "headers": {
                                    "type": "array",
                                    "description": "表头行（可选），会加粗并设置底色",
                                    "items": {"type": "string"}
                                },
                                "rows": {
                                    "type": "array",
                                    "description": "数据行，每行是一个单元格值数组（字符串/数字/布尔）",
                                    "items": {
                                        "type": "array",
                                        "items": {}
                                    }
                                },
                                "start_cell": {
                                    "type": "string",
                                    "description": "数据起始单元格，如 'A1'（默认）或 'B2'",
                                    "default": "A1"
                                }
                            },
                            "required": ["rows"],
                            "additionalProperties": False
                        }
                    },
                    "overwrite": {
                        "type": "boolean",
                        "description": "若文件已存在是否覆盖。默认 true；设为 false 时改为追加新工作表到已有文件。",
                        "default": True
                    }
                },
                "required": ["file_path", "sheets"],
                "additionalProperties": False
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "write_pptx",
            "description": (
                "创建 PowerPoint 演示文稿 (.pptx)。"
                "可指定演示文稿标题（生成封面页），每张幻灯片可包含标题、"
                "要点列表（支持缩进层级）和演讲者备注。"
                "适合生成汇报、培训、提案、路演等演示材料。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": (
                            "输出文件路径。例如 'slides/pitch.pptx'。"
                            "扩展名必须为 .pptx，父目录不存在时会自动创建。"
                        )
                    },
                    "slides": {
                        "type": "array",
                        "description": (
                            "幻灯片列表，每个元素格式：\n"
                            "{\"title\": \"标题\", \"bullets\": [\"要点1\", {\"text\": \"子要点\", \"level\": 1}], "
                            "\"notes\": \"演讲者备注\"}\n"
                            "bullets 中字符串为一级要点（圆点），"
                            "对象可指定 level（0-4）控制缩进层级。"
                        ),
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "幻灯片标题"},
                                "bullets": {
                                    "type": "array",
                                    "description": "要点列表，支持字符串或 {\"text\": ..., \"level\": n}",
                                    "items": {}
                                },
                                "notes": {"type": "string", "description": "演讲者备注（可选）"}
                            },
                            "required": ["title"],
                            "additionalProperties": False
                        }
                    },
                    "title": {
                        "type": "string",
                        "description": "演示文稿总标题（可选），会生成一张封面页。",
                        "default": ""
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "封面页副标题（可选，仅当提供 title 时生效）。",
                        "default": ""
                    },
                    "include_notes": {
                        "type": "boolean",
                        "description": "是否写入演讲者备注，默认 true。",
                        "default": True
                    },
                    "overwrite": {
                        "type": "boolean",
                        "description": "若文件已存在是否覆盖，默认 true。",
                        "default": True
                    }
                },
                "required": ["file_path", "slides"],
                "additionalProperties": False
            }
        }
    }
]

# ═══════════════════════════════════════════════════════════════
#  Lazy imports – 优雅处理依赖缺失
# ═══════════════════════════════════════════════════════════════

_MISSING_DOCX: Optional[str] = None
_MISSING_XLSX: Optional[str] = None
_MISSING_PPTX: Optional[str] = None


def _get_docx_module():
    """Lazy-import python-docx with a helpful error message."""
    global _MISSING_DOCX
    if _MISSING_DOCX:
        return None, _MISSING_DOCX
    try:
        import docx
        return docx, None
    except ImportError:
        _MISSING_DOCX = (
            "❌ 未安装 python-docx 库。\n"
            "请运行：`pip install python-docx`\n"
            "然后重试。"
        )
        return None, _MISSING_DOCX


def _get_xlsx_module():
    """Lazy-import openpyxl with a helpful error message."""
    global _MISSING_XLSX
    if _MISSING_XLSX:
        return None, _MISSING_XLSX
    try:
        import openpyxl
        return openpyxl, None
    except ImportError:
        _MISSING_XLSX = (
            "❌ 未安装 openpyxl 库。\n"
            "请运行：`pip install openpyxl`\n"
            "然后重试。"
        )
        return None, _MISSING_XLSX


def _get_pptx_module():
    """Lazy-import python-pptx with a helpful error message."""
    global _MISSING_PPTX
    if _MISSING_PPTX:
        return None, _MISSING_PPTX
    try:
        import pptx
        return pptx, None
    except ImportError:
        _MISSING_PPTX = (
            "❌ 未安装 python-pptx 库。\n"
            "请运行：`pip install python-pptx`\n"
            "然后重试。"
        )
        return None, _MISSING_PPTX

# ═══════════════════════════════════════════════════════════════
#  通用 Helpers
# ═══════════════════════════════════════════════════════════════


def _resolve_path(file_path: str, context) -> str:
    """Resolve file path relative to workspace root or as absolute."""
    if os.path.isabs(file_path):
        return file_path
    project_root = getattr(context, "project_root", "") or os.getcwd()
    return os.path.normpath(os.path.join(project_root, file_path))


def _validate_extension(file_path: str, expected: str) -> Optional[str]:
    """Check the file extension matches the expected one."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext != expected:
        return (
            f"❌ 文件扩展名不匹配：期望 `{expected}`，"
            f"实际为 `{ext or '(无扩展名)'}`"
        )
    return None


def _ensure_parent_dir(file_path: str):
    """Create parent directories if they don't exist."""
    parent = os.path.dirname(file_path)
    if parent and not os.path.exists(parent):
        os.makedirs(parent, exist_ok=True)


def _coerce_list(value, param_name: str):
    """Accept a JSON string or a list, and return a list.

    Returns (list, error_msg) – error_msg is non-empty on failure.
    """
    if value is None:
        return [], ""
    if isinstance(value, list):
        return value, ""
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
            if isinstance(parsed, list):
                return parsed, ""
            return [], f"❌ 参数 `{param_name}` 解析后不是数组。"
        except json.JSONDecodeError as exc:
            return [], f"❌ 参数 `{param_name}` 不是合法的 JSON 数组：{exc}"
    return [], f"❌ 参数 `{param_name}` 必须是数组或 JSON 字符串。"


def _fmt_size(size: int) -> str:
    """Format byte count in a human-friendly way."""
    if size < 1024:
        return f"{size} 字节"
    if size < 1024 * 1024:
        return f"{size / 1024:.1f} KB"
    return f"{size / 1024 / 1024:.2f} MB"

# ═══════════════════════════════════════════════════════════════
#  DOCX Writer (Markdown → Word)
# ═══════════════════════════════════════════════════════════════

# 行内格式：**粗体**、*斜体*、`代码`（** 必须优先匹配）
_INLINE_RE = re.compile(r"(\*\*.+?\*\*|`[^`]+`|\*[^*]+?\*)")


def _add_runs(paragraph, text: str):
    """Parse inline markdown (**bold**, *italic*, `code`) into runs."""
    pos = 0
    for match in _INLINE_RE.finditer(text):
        if match.start() > pos:
            paragraph.add_run(text[pos:match.start()])
        token = match.group(0)
        if token.startswith("**"):
            run = paragraph.add_run(token[2:-2])
            run.bold = True
        elif token.startswith("`"):
            run = paragraph.add_run(token[1:-1])
            run.font.name = "Consolas"
            try:
                from docx.shared import Pt
                run.font.size = Pt(10)
            except ImportError:
                pass
        else:
            run = paragraph.add_run(token[1:-1])
            run.italic = True
        pos = match.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


def _add_horizontal_rule(doc):
    """Add a horizontal rule via paragraph bottom border (OXML)."""
    try:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "999999")
        pBdr.append(bottom)
        pPr.append(pBdr)
    except Exception:
        # Fallback: a subtle divider paragraph
        doc.add_paragraph("─" * 30)


def _split_table_row(line: str) -> List[str]:
    """Split a markdown table row '| a | b |' into cells."""
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [cell.strip() for cell in line.split("|")]


def _is_separator_row(cells: List[str]) -> bool:
    """Detect '|---|---|' style separator rows."""
    return bool(cells) and all(
        re.fullmatch(r":?-{2,}:?", c) for c in cells
    )


def _collect_table_lines(lines: List[str], idx: int) -> Tuple[List[str], int]:
    """Collect consecutive markdown table rows starting at idx."""
    table_lines: List[str] = []
    while idx < len(lines) and lines[idx].strip().startswith("|"):
        table_lines.append(lines[idx].strip())
        idx += 1
    return table_lines, idx


def _add_md_table(doc, table_lines: List[str], stats: Dict[str, int]):
    """Convert markdown table lines into a docx table."""
    parsed = [_split_table_row(l) for l in table_lines]
    parsed = [r for r in parsed if not _is_separator_row(r)]
    if not parsed:
        return

    headers = parsed[0]
    body = parsed[1:]
    n_cols = len(headers)
    n_rows = len(body) + 1  # 表头 + 数据行

    table = doc.add_table(rows=n_rows, cols=n_cols)
    try:
        table.style = "Light Grid Accent 1"
    except KeyError:
        try:
            table.style = "Table Grid"
        except KeyError:
            pass

    # 表头（加粗）
    for j, h in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = ""
        run = cell.paragraphs[0].add_run(h)
        run.bold = True

    # 数据行
    for i, row in enumerate(body, start=1):
        for j in range(n_cols):
            val = row[j] if j < len(row) else ""
            table.rows[i].cells[j].text = val

    stats["tables"] += 1
    stats["paragraphs"] += 1


def _md_to_docx(doc, md_text: str) -> Dict[str, int]:
    """Convert Markdown text into docx content. Returns stats dict."""
    stats = {"headings": 0, "paragraphs": 0, "tables": 0, "lists": 0}
    lines = md_text.split("\n")
    idx = 0
    n = len(lines)

    while idx < n:
        line = lines[idx].rstrip()
        stripped = line.strip()

        # ── 空行 ──
        if not stripped:
            idx += 1
            continue

        # ── 表格 ──
        if stripped.startswith("|"):
            table_lines, idx = _collect_table_lines(lines, idx)
            _add_md_table(doc, table_lines, stats)
            continue

        # ── 分隔线 --- / *** / ___ ──
        if re.fullmatch(r"([-*_])\1{2,}", stripped):
            _add_horizontal_rule(doc)
            idx += 1
            continue

        # ── 标题 ──
        heading = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading:
            level = len(heading.group(1))
            doc.add_heading(heading.group(2).strip(), level=level)
            stats["headings"] += 1
            stats["paragraphs"] += 1
            idx += 1
            continue

        # ── 引用（连续 > 行合并为一段）──
        if stripped.startswith(">"):
            quote_parts: List[str] = []
            while idx < n and lines[idx].strip().startswith(">"):
                quote_parts.append(lines[idx].strip().lstrip(">").strip())
                idx += 1
            p = doc.add_paragraph()
            try:
                p.style = doc.styles["Intense Quote"]
            except KeyError:
                pass
            _add_runs(p, " ".join(quote_parts))
            stats["paragraphs"] += 1
            continue

        # ── 列表（缩进 2 空格 = 1 级）──
        list_match = re.match(r"^(\s*)([-*+]|\d+[.)])\s+(.*)$", line)
        if list_match:
            indent = len(list_match.group(1))
            marker = list_match.group(2)
            text = list_match.group(3)
            ordered = marker[0].isdigit()
            level = min(indent // 2, 3)

            style_name = f"List {'Number' if ordered else 'Bullet'}"
            if level > 0:
                style_name += f" {level + 1}"

            p = doc.add_paragraph()
            try:
                p.style = doc.styles[style_name]
            except KeyError:
                try:
                    p.style = doc.styles["List Bullet" if not ordered else "List Number"]
                except KeyError:
                    pass
            _add_runs(p, text)
            stats["lists"] += 1
            stats["paragraphs"] += 1
            idx += 1
            continue

        # ── 普通段落 ──
        p = doc.add_paragraph()
        _add_runs(p, stripped)
        stats["paragraphs"] += 1
        idx += 1

    return stats


def _write_docx(file_path: str, content: str, title: str = "",
                mode: str = "create", overwrite: bool = True) -> str:
    """Create or append a Word document from Markdown content."""
    docx, err = _get_docx_module()
    if err:
        return err

    ext_err = _validate_extension(file_path, ".docx")
    if ext_err:
        return ext_err

    exists = os.path.exists(file_path)

    if mode == "append":
        if exists:
            try:
                doc = docx.Document(file_path)
            except Exception as exc:
                return f"❌ 无法打开现有文档 `{file_path}`：{exc}"
        else:
            doc = docx.Document()
            if title:
                doc.add_heading(title, level=0)
    else:
        if exists and not overwrite:
            return (
                f"❌ 文件已存在：`{file_path}`。\n"
                f"如需覆盖请设置 overwrite=True，或改用 mode=\"append\" 追加内容。"
            )
        doc = docx.Document()
        if title:
            doc.add_heading(title, level=0)

    stats = _md_to_docx(doc, content)

    try:
        _ensure_parent_dir(file_path)
        doc.save(file_path)
    except Exception as exc:
        return f"❌ 保存文档失败：{exc}"

    size = os.path.getsize(file_path)
    action = "追加写入" if (mode == "append" and exists) else "写入"
    return (
        f"✅ 已{action} Word 文档 `{file_path}`（{_fmt_size(size)}）。\n"
        f"📊 内容统计：{stats['headings']} 个标题，{stats['paragraphs']} 个段落/块，"
        f"{stats['tables']} 个表格，{stats['lists']} 个列表项。\n"
        f"💡 可调用 `read_docx` 工具读取验证。"
    )

# ═══════════════════════════════════════════════════════════════
#  XLSX Writer (结构化数据 → Excel)
# ═══════════════════════════════════════════════════════════════


def _parse_cell_ref(ref: str) -> Tuple[int, int]:
    """Parse 'B2' → (col=2, row=2). Defaults to (1, 1)."""
    match = re.match(r"^([A-Za-z]+)(\d+)$", (ref or "").strip())
    if not match:
        return (1, 1)
    col = 0
    for ch in match.group(1).upper():
        col = col * 26 + (ord(ch) - ord("A") + 1)
    return (col, int(match.group(2)))


def _autofit_columns(ws, min_width: int = 8, max_width: int = 50):
    """Auto-size column widths (CJK chars count double)."""
    try:
        from openpyxl.utils import get_column_letter
    except ImportError:
        return
    for col_cells in ws.columns:
        if not col_cells:
            continue
        col_idx = col_cells[0].column
        max_len = 0
        for cell in col_cells:
            if cell.value is None:
                continue
            s = str(cell.value)
            length = sum(2 if ord(ch) > 127 else 1 for ch in s)
            max_len = max(max_len, length)
        width = min(max(max_len + 2, min_width), max_width)
        ws.column_dimensions[get_column_letter(col_idx)].width = width


def _write_xlsx(file_path: str, sheets: List[dict],
                overwrite: bool = True) -> str:
    """Create (or append to) an Excel workbook with multiple sheets."""
    openpyxl, err = _get_xlsx_module()
    if err:
        return err

    ext_err = _validate_extension(file_path, ".xlsx")
    if ext_err:
        return ext_err

    exists = os.path.exists(file_path)

    if exists and not overwrite:
        # 追加模式：打开已有工作簿
        try:
            wb = openpyxl.load_workbook(file_path)
        except Exception as exc:
            return f"❌ 无法打开现有工作簿 `{file_path}`：{exc}"
        base_msg = "已追加工作表到"
    else:
        if exists and overwrite:
            base_msg = "已覆盖并写入"
        else:
            base_msg = "已写入"
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # 移除默认空表，按需创建

    added: List[Tuple[str, int, int]] = []
    used_names = set(wb.sheetnames)

    for spec in sheets:
        if not isinstance(spec, dict):
            return f"❌ 工作表定义必须是对象，收到：{type(spec).__name__}"

        name = str(spec.get("name") or f"Sheet{len(added) + 1}")
        headers = spec.get("headers") or []
        rows = spec.get("rows") or []
        start_cell = str(spec.get("start_cell") or "A1")

        # 工作表名去重
        base_name, n = name, 1
        while base_name in used_names:
            n += 1
            base_name = f"{name} ({n})"
        used_names.add(base_name)

        ws = wb.create_sheet(title=base_name)
        col_start, row_start = _parse_cell_ref(start_cell)

        # 表头
        if headers:
            for j, h in enumerate(headers):
                cell = ws.cell(row=row_start, column=col_start + j, value=h)
                cell.font = openpyxl.styles.Font(bold=True)
                cell.fill = openpyxl.styles.PatternFill(
                    "solid", fgColor="D9E1F2")
            row_start += 1

        # 数据
        for row in rows:
            if not isinstance(row, (list, tuple)):
                row = [row]
            for j, val in enumerate(row):
                if val is not None:
                    ws.cell(row=row_start, column=col_start + j, value=val)
            row_start += 1

        # 冻结表头 + 自动列宽
        if headers:
            ws.freeze_panes = ws.cell(row=row_start, column=col_start)
        _autofit_columns(ws)

        added.append((base_name, max(len(rows), 1), len(headers) or (len(rows[0]) if rows else 0)))

    if not added:
        return "❌ sheets 参数为空，请至少提供一个工作表定义。"

    try:
        _ensure_parent_dir(file_path)
        wb.save(file_path)
    except Exception as exc:
        return f"❌ 保存工作簿失败：{exc}"

    size = os.path.getsize(file_path)
    lines = [f"✅ {base_msg} Excel 工作簿 `{file_path}`（{_fmt_size(size)}）。"]
    for name, r, c in added:
        lines.append(f"📊 `{name}`：{r} 行 × {c} 列")
    lines.append("💡 可调用 `read_xlsx` 工具读取验证。")
    return "\n".join(lines)

# ═══════════════════════════════════════════════════════════════
#  PPTX Writer (结构化数据 → PowerPoint)
# ═══════════════════════════════════════════════════════════════


def _get_placeholder(slide, idx: int):
    """Find a placeholder by its format index (e.g. 0=title, 1=body)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _write_pptx(file_path: str, slides: List[dict], title: str = "",
                subtitle: str = "", include_notes: bool = True,
                overwrite: bool = True) -> str:
    """Create a PowerPoint presentation from structured slide data."""
    pptx, err = _get_pptx_module()
    if err:
        return err

    ext_err = _validate_extension(file_path, ".pptx")
    if ext_err:
        return ext_err

    exists = os.path.exists(file_path)
    if exists and not overwrite:
        return (
            f"❌ 文件已存在：`{file_path}`。\n"
            f"如需覆盖请设置 overwrite=True。"
        )

    if not slides and not title:
        return "❌ slides 为空且未提供 title，无法生成空的演示文稿。"

    prs = pptx.Presentation()
    slide_count = 0

    # ── 封面页 ──
    if title:
        cover = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        cover.shapes.title.text = title
        if subtitle:
            sub_ph = _get_placeholder(cover, 1)
            if sub_ph is not None:
                sub_ph.text = subtitle
        slide_count += 1

    # ── 内容页 ──
    for spec in slides:
        if not isinstance(spec, dict):
            return f"❌ 幻灯片定义必须是对象，收到：{type(spec).__name__}"

        s_title = str(spec.get("title", ""))
        bullets = spec.get("bullets") or []
        notes = str(spec.get("notes", ""))

        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
        if s_title:
            slide.shapes.title.text = s_title

        body = _get_placeholder(slide, 1)
        if body is not None:
            tf = body.text_frame
            tf.clear()
            first = True
            for item in bullets:
                if isinstance(item, dict):
                    text = str(item.get("text", ""))
                    try:
                        level = min(max(int(item.get("level", 0)), 0), 4)
                    except (TypeError, ValueError):
                        level = 0
                else:
                    text = str(item)
                    level = 0
                if not text:
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = text
                p.level = level

        if notes and include_notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass  # 备注失败不影响主内容

        slide_count += 1

    try:
        _ensure_parent_dir(file_path)
        prs.save(file_path)
    except Exception as exc:
        return f"❌ 保存演示文稿失败：{exc}"

    size = os.path.getsize(file_path)
    cover_note = "（含 1 张封面页）" if title else ""
    return (
        f"✅ 已写入 PowerPoint 演示文稿 `{file_path}`（{_fmt_size(size)}）。\n"
        f"📽️ 共 {slide_count} 张幻灯片{cover_note}。\n"
        f"💡 可调用 `read_pptx` 工具读取验证。"
    )

# ═══════════════════════════════════════════════════════════════
#  Hook: 启动时依赖检查
# ═══════════════════════════════════════════════════════════════


def on_agent_init(context):
    """Check dependency availability at startup."""
    status: List[str] = []

    checks = (
        ("DOCX", _get_docx_module),
        ("XLSX", _get_xlsx_module),
        ("PPTX", _get_pptx_module),
    )
    deps = {}
    for label, fn in checks:
        mod, _ = fn()
        deps[label.lower()] = mod is not None
        status.append(f"✅ {label}" if mod else f"⚠️ {label}（需安装依赖）")

    context.logger.info(f"Office Writer loaded — {', '.join(status)}")
    context.storage["office_writer_deps"] = deps

# ═══════════════════════════════════════════════════════════════
#  工具分发
# ═══════════════════════════════════════════════════════════════


def execute(tool_name: str, args: dict, context) -> str:
    """Dispatch tool calls to the appropriate handler."""

    if tool_name == "write_docx":
        file_path = _resolve_path(args.get("file_path", ""), context)
        if not file_path:
            return "❌ 请提供 file_path 参数。"
        content = args.get("content", "")
        if not content or not content.strip():
            return "❌ content 参数为空，请提供要写入的 Markdown 内容。"
        return _write_docx(
            file_path=file_path,
            content=content,
            title=args.get("title", ""),
            mode=args.get("mode", "create"),
            overwrite=args.get("overwrite", True),
        )

    if tool_name == "write_xlsx":
        file_path = _resolve_path(args.get("file_path", ""), context)
        if not file_path:
            return "❌ 请提供 file_path 参数。"
        sheets, err = _coerce_list(args.get("sheets"), "sheets")
        if err:
            return err
        if not sheets:
            return "❌ sheets 参数为空，请至少提供一个工作表（含 rows 数据）。"
        return _write_xlsx(
            file_path=file_path,
            sheets=sheets,
            overwrite=args.get("overwrite", True),
        )

    if tool_name == "write_pptx":
        file_path = _resolve_path(args.get("file_path", ""), context)
        if not file_path:
            return "❌ 请提供 file_path 参数。"
        slides, err = _coerce_list(args.get("slides"), "slides")
        if err:
            return err
        if not slides and not args.get("title"):
            return "❌ slides 为空且未提供 title，无法生成空的演示文稿。"
        return _write_pptx(
            file_path=file_path,
            slides=slides,
            title=args.get("title", ""),
            subtitle=args.get("subtitle", ""),
            include_notes=args.get("include_notes", True),
            overwrite=args.get("overwrite", True),
        )

    return f"Unknown tool: {tool_name}"
